from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterator

from .pdf_layout import (
    PDF_LAYOUT_ADAPTER,
    SourceSegment,
    layout_sidecar,
    parse_pdf_layout,
    semantic_eligible_text,
)


class ParseError(RuntimeError):
    pass


@dataclass
class ParsedSource:
    text: str
    source_type: str
    diagnostics: dict[str, Any]
    segments: tuple[SourceSegment, ...] | None = None
    layout_sidecar: dict[str, Any] | None = None


@dataclass
class _ParseStats:
    total_units: int = 0
    text_units: int = 0
    error_units: int = 0
    extracted_chars: int = 0

    def record(self, text: str, *, error: bool = False) -> None:
        self.total_units += 1
        self.error_units += int(error)
        chars = sum(not char.isspace() for char in text)
        self.text_units += int(chars > 0)
        self.extracted_chars += chars


# Parser, locator scheme and counting unit. Format remains the original extension.
FORMAT_DETAILS = {
    **{ext: ("builtin.text", "TEXT", "document") for ext in ("txt", "md", "markdown", "csv")},
    "pdf": ("pypdf", "PAGE", "page"),
    "docx": ("python-docx", "PARA / TABLE", "paragraph_or_table_row"),
    "xlsx": ("openpyxl", "SHEET / ROW", "row"),
    "xlsm": ("openpyxl", "SHEET / ROW", "row"),
    "pptx": ("python-pptx", "SLIDE", "slide"),
}
LOCATOR_PATTERN = r"(?:PAGE:[1-9]\d*|PARA:[1-9]\d*|TABLE:[1-9]\d*(?::ROW:[1-9]\d*)?|SHEET:[^\\/*?:\[\]\r\n]{1,31}(?::ROW:[1-9]\d*)?|SLIDE:[1-9]\d*)"
SOURCE_MARKER = re.compile(r"^\[\[(" + LOCATOR_PATTERN + r")\]\]", re.MULTILINE)


def source_units(text: str) -> Iterator[tuple[str, str]]:
    """Yield marker-free units in Source order; unmarked text has no page number."""
    markers = list(SOURCE_MARKER.finditer(text))
    prefix_end = markers[0].start() if markers else len(text)
    if text[:prefix_end].strip():
        yield "TEXT", text[:prefix_end]
    for index, marker in enumerate(markers):
        end = markers[index + 1].start() if index + 1 < len(markers) else len(text)
        body = text[marker.end():end]
        if body.strip() != "[PAGE_PARSE_ERROR]":
            yield marker.group(1), body


def parse_warnings(diagnostics: dict[str, Any]) -> list[str]:
    warnings = []
    if diagnostics.get("partial_parse"):
        warnings.append("Partial extraction; some units could not be parsed.")
    if diagnostics.get("error_units"):
        warnings.append(f"{diagnostics['error_units']} {diagnostics['unit_type']} parse errors.")
    if diagnostics.get("empty_units"):
        warnings.append(f"{diagnostics['empty_units']} empty {diagnostics['unit_type']} units.")
    if diagnostics.get("empty_extraction"):
        warnings.append(
            "No extractable text; OCR/multimodal parsing required."
            if diagnostics.get("image_only_or_no_extractable_text")
            else "No extractable text."
        )
    return warnings


def parse_text_file(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "gb18030"):
        try:
            return path.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    raise ParseError(f"Unable to decode text file: {path}")


def parse_pdf(path: Path, *, _stats: _ParseStats | None = None) -> str:
    from pypdf import PdfReader
    stats = _stats or _ParseStats()
    reader = PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception:
            stats.record("", error=True)
            text = "[PAGE_PARSE_ERROR]"
        else:
            stats.record(text)
        parts.append(f"\n[[PAGE:{i}]]\n{text}")
    return "\n".join(parts)


def parse_docx(path: Path, *, _stats: _ParseStats | None = None) -> str:
    from docx import Document
    stats = _stats or _ParseStats()
    doc = Document(str(path))
    parts = []
    for i, p in enumerate(doc.paragraphs, 1):
        stats.record(p.text)
        if p.text.strip():
            parts.append(f"[[PARA:{i}]] {p.text}")
    for ti, table in enumerate(doc.tables, 1):
        parts.append(f"[[TABLE:{ti}]]")
        for ri, row in enumerate(table.rows, 1):
            vals = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            stats.record("".join(vals))
            parts.append(f"[[TABLE:{ti}:ROW:{ri}]] " + " | ".join(vals))
    return "\n".join(parts)


def parse_xlsx(path: Path, *, _stats: _ParseStats | None = None) -> str:
    from openpyxl import load_workbook
    stats = _stats or _ParseStats()
    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts = []
    try:
        for ws in wb.worksheets:
            parts.append(f"[[SHEET:{ws.title}]]")
            for ri, row in enumerate(ws.iter_rows(values_only=True), 1):
                vals = ["" if v is None else str(v) for v in row]
                stats.record("".join(vals))
                if any(v.strip() for v in vals):
                    parts.append(f"[[SHEET:{ws.title}:ROW:{ri}]] " + " | ".join(vals))
    finally:
        wb.close()
    return "\n".join(parts)


def parse_pptx(path: Path, *, _stats: _ParseStats | None = None) -> str:
    from pptx import Presentation
    stats = _stats or _ParseStats()
    prs = Presentation(str(path))
    parts = []
    for si, slide in enumerate(prs.slides, 1):
        parts.append(f"[[SLIDE:{si}]]")
        slide_text = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
                slide_text.append(text.strip())
        stats.record("".join(slide_text))
    return "\n".join(parts)


def parse_source(path: Path) -> tuple[str, str]:
    parsed = parse_source_with_diagnostics(path)
    return parsed.text, parsed.source_type


def parse_source_with_diagnostics(
    path: Path, *, include_semantic_segments: bool = False,
) -> ParsedSource:
    path = Path(path)
    fmt = path.suffix.lower().lstrip(".")
    stats = _ParseStats()
    if fmt in {"txt", "md", "markdown", "csv"}:
        text = parse_text_file(path)
        stats.record(text)
    elif fmt in {"pdf", "docx", "xlsx", "xlsm", "pptx"}:
        parser = {"pdf": parse_pdf, "docx": parse_docx, "xlsx": parse_xlsx,
                  "xlsm": parse_xlsx, "pptx": parse_pptx}[fmt]
        text = parser(path, _stats=stats)
    elif fmt in {"png", "jpg", "jpeg", "webp"}:
        raise ParseError("Image OCR/multimodal parsing is intentionally deferred in v0.1; file can still be archived/uploaded.")
    else:
        raise ParseError(f"Unsupported parser for extension: {path.suffix.lower() or '<none>'}")
    parser_name, scheme, unit_type = FORMAT_DETAILS[fmt]
    diagnostics = {
        "format": fmt, "parser": parser_name, "locator_scheme": scheme,
        "file_size": path.stat().st_size, "unit_type": unit_type,
        "total_units": stats.total_units, "text_units": stats.text_units,
        "error_units": stats.error_units,
        "empty_units": stats.total_units - stats.text_units - stats.error_units,
        "extracted_chars": stats.extracted_chars,
        "empty_extraction": stats.extracted_chars == 0,
        "partial_parse": stats.error_units > 0 and stats.text_units > 0,
        "image_only_or_no_extractable_text": fmt == "pdf" and stats.total_units > 0 and stats.extracted_chars == 0,
    }
    segments = None
    sidecar = None
    if fmt == "pdf" and include_semantic_segments:
        projection = parse_pdf_layout(path, text)
        segments = projection.segments
        sidecar = layout_sidecar(projection)
        diagnostics["pdf_layout"] = {
            "adapter": PDF_LAYOUT_ADAPTER,
            "adapter_versions": dict(projection.adapter_versions),
            "signature_sha256": projection.signature_sha256,
            "segments": len(segments),
            "segment_kind_counts": {
                kind: sum(segment.kind == kind for segment in segments)
                for kind in ("narrative", "table", "unknown")
            },
            "semantic_policy": "NARRATIVE_FIRST_TABLE_SUPPRESSION",
            "canonical_text_parser_unchanged": True,
        }
    return ParsedSource(
        text, "xlsx" if fmt == "xlsm" else fmt, diagnostics, segments, sidecar,
    )


def semantic_eligible_source_text(parsed: ParsedSource) -> str:
    """Return the pre-chunk semantic view without mutating Source truth."""
    return semantic_eligible_text(parsed.text, parsed.segments)


def chunk_text(text: str, max_chars: int) -> list[str]:
    if len(text) <= max_chars:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + max_chars)
        if end < len(text):
            boundary = text.rfind("\n", start, end)
            if boundary > start + max_chars // 2:
                end = boundary
        chunks.append(text[start:end])
        start = end
    return chunks


def chunk_source_text(text: str, max_chars: int) -> list[str]:
    """Pack locator units, splitting oversized units without duplicating any text."""
    if max_chars <= 0:
        raise ValueError("max_chars must be positive")
    markers = list(SOURCE_MARKER.finditer(text))
    if not markers:
        return chunk_text(text, max_chars)
    if any(len(marker.group()) > max_chars for marker in markers):
        raise ValueError("max_chars is smaller than a locator marker")
    boundaries = sorted({0, *(marker.start() for marker in markers), len(text)})
    first = markers[0]
    if first.start() and not text[:first.start()].strip() and first.end() < max_chars:
        # PDF's leading separator belongs with its first page, not an empty prompt.
        boundaries.remove(first.start())
    chunks: list[str] = []
    pending = ""
    for start, end in zip(boundaries, boundaries[1:]):
        unit = text[start:end]
        if len(pending) + len(unit) <= max_chars:
            pending += unit
            continue
        if pending:
            chunks.append(pending)
            pending = ""
        if len(unit) <= max_chars:
            pending = unit
            continue
        marker = SOURCE_MARKER.search(unit)
        # Do not choose the newline immediately after a marker as a soft boundary.
        body_start = len(unit) - len(unit[marker.end():].lstrip()) if marker else 0
        offset = 0
        while offset < len(unit):
            cut = min(len(unit), offset + max_chars)
            if cut < len(unit):
                newline = unit.rfind("\n", offset, cut)
                if newline > max(offset + max_chars // 2, body_start):
                    cut = newline
            chunks.append(unit[offset:cut])
            offset = cut
    if pending:
        chunks.append(pending)
    return chunks
