"""Tiny generated Sources; no user files, network or committed binary fixtures."""
from pathlib import Path


EXCERPT = "Capacity reached 42 units."
FORMATS = ("txt", "md", "markdown", "csv", "pdf", "docx", "xlsx", "xlsm", "pptx")


def write_pdf(path: Path, pages: list[str]) -> None:
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    for text in pages:
        page = writer.add_blank_page(width=300, height=200)
        if not text:
            continue
        font = DictionaryObject({NameObject("/Type"): NameObject("/Font"),
                                 NameObject("/Subtype"): NameObject("/Type1"),
                                 NameObject("/BaseFont"): NameObject("/Helvetica")})
        page[NameObject("/Resources")] = DictionaryObject({
            NameObject("/Font"): DictionaryObject({NameObject("/F1"): font}),
        })
        stream = DecodedStreamObject()
        stream.set_data(f"BT /F1 12 Tf 10 100 Td ({text}) Tj ET".encode("ascii"))
        # An indirect stream gives real pypdf extraction coverage without reportlab.
        page[NameObject("/Contents")] = writer._add_object(stream)
    writer.write(path)


def write_source(path: Path) -> str:
    fmt = path.suffix.lower().lstrip(".")
    if fmt in {"txt", "md", "markdown", "csv"}:
        path.write_text(f"Introduction\n{EXCERPT}\nClosing note", encoding="utf-8")
        return "TEXT"
    if fmt == "pdf":
        write_pdf(path, ["Introduction", EXCERPT, "Closing note"])
        return "PAGE:2"
    if fmt == "docx":
        from docx import Document
        doc = Document()
        doc.add_paragraph("Introduction")
        doc.add_paragraph("")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Capacity"
        table.cell(1, 0).text = EXCERPT
        doc.save(path)
        return "TABLE:1:ROW:2"
    if fmt in {"xlsx", "xlsm"}:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Capacity"
        ws["A1"] = "Introduction"
        ws["A3"] = EXCERPT
        wb.create_sheet("Notes")["A1"] = "Closing note"
        wb.save(path)
        wb.close()
        return "SHEET:Capacity:ROW:3"
    if fmt == "pptx":
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        for lines in (("Introduction", "First slide"), (EXCERPT, "Closing note")):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            for index, text in enumerate(lines):
                slide.shapes.add_textbox(Inches(1), Inches(index + 1), Inches(6), Inches(1)).text = text
        prs.save(path)
        return "SLIDE:2"
    raise AssertionError(fmt)
