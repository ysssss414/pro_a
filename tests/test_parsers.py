from types import SimpleNamespace
from unittest.mock import Mock

import pytest

from pro_a.analyzer import evidence_match, resolve_evidence_locator
from pro_a.parsers import (
    FORMAT_DETAILS, SOURCE_MARKER, ParseError, chunk_source_text, chunk_text,
    parse_source, parse_source_with_diagnostics, parse_warnings,
)

from multiformat_helpers import EXCERPT, FORMATS, write_source


@pytest.mark.parametrize("fmt", FORMATS)
def test_multiformat_diagnostics_locators_and_compatibility(tmp_path, fmt):
    path = tmp_path / f"fixture.{fmt}"
    locator = write_source(path)
    parsed = parse_source_with_diagnostics(path)
    assert parsed.source_type == ("xlsx" if fmt == "xlsm" else fmt)
    assert parse_source(path) == (parsed.text, parsed.source_type)
    diag = parsed.diagnostics
    assert diag == parse_source_with_diagnostics(path).diagnostics
    assert (diag["parser"], diag["locator_scheme"], diag["unit_type"]) == FORMAT_DETAILS[fmt]
    assert diag["format"] == fmt
    assert diag["file_size"] == path.stat().st_size
    assert diag["total_units"] == diag["text_units"] + diag["empty_units"] + diag["error_units"]
    assert diag["text_units"] > 0 and diag["extracted_chars"] > 0
    assert diag["empty_extraction"] is False and diag["partial_parse"] is False
    if locator != "TEXT":
        assert f"[[{locator}]]" in parsed.text
    else:
        assert not list(SOURCE_MARKER.finditer(parsed.text))
    assert resolve_evidence_locator(parsed.text, EXCERPT) == {"status": "resolved", "locator": locator}
    chunks = chunk_source_text(parsed.text, 70)
    assert "".join(chunks) == parsed.text
    assert max(map(len, chunks)) <= 70
    for marker in SOURCE_MARKER.finditer(parsed.text):
        assert sum(chunk.count(marker.group()) for chunk in chunks) == 1


@pytest.mark.parametrize("fmt,markers,total,text_units", [
    ("pdf", ["PAGE:1", "PAGE:2", "PAGE:3"], 3, 3),
    ("docx", ["PARA:1", "TABLE:1", "TABLE:1:ROW:1", "TABLE:1:ROW:2"], 4, 3),
    ("xlsx", ["SHEET:Capacity", "SHEET:Capacity:ROW:3", "SHEET:Notes:ROW:1"], 4, 3),
    ("pptx", ["SLIDE:1", "SLIDE:2"], 2, 2),
])
def test_unit_counts_and_existing_markers(tmp_path, fmt, markers, total, text_units):
    path = tmp_path / f"fixture.{fmt}"
    write_source(path)
    parsed = parse_source_with_diagnostics(path)
    for marker in markers:
        assert f"[[{marker}]]" in parsed.text
    assert parsed.diagnostics["total_units"] == total
    assert parsed.diagnostics["text_units"] == text_units
    if fmt == "pptx":
        assert "First slide" in parsed.text and "Closing note" in parsed.text


@pytest.mark.parametrize("values,text_units,errors,empty", [
    (["alpha", "beta", "gamma"], 3, 0, False),
    (["alpha", RuntimeError("private path must not escape"), "gamma"], 2, 1, False),
    ([None, "", " \n\t"], 0, 0, True),
    ([RuntimeError("bad"), RuntimeError("bad")], 0, 2, True),
    ([], 0, 0, True),
])
def test_pdf_page_quality(tmp_path, monkeypatch, values, text_units, errors, empty):
    path = tmp_path / "fixture.pdf"
    path.write_bytes(b"mock PDF")
    pages = [SimpleNamespace(extract_text=Mock(side_effect=value) if isinstance(value, Exception)
                             else Mock(return_value=value)) for value in values]
    monkeypatch.setattr("pypdf.PdfReader", lambda _: SimpleNamespace(pages=pages))
    parsed = parse_source_with_diagnostics(path)
    diag = parsed.diagnostics
    assert diag["total_units"] == len(values)
    assert diag["text_units"] == text_units and diag["error_units"] == errors
    assert diag["empty_extraction"] is empty
    assert diag["partial_parse"] is bool(errors and text_units)
    assert diag["image_only_or_no_extractable_text"] is bool(values and empty)
    assert diag["extracted_chars"] == sum(sum(not c.isspace() for c in value)
                                          for value in values if isinstance(value, str))
    assert parsed.text.count("[PAGE_PARSE_ERROR]") == errors
    assert "private path" not in parsed.text
    if errors and text_units:
        assert "Partial extraction" in " ".join(parse_warnings(diag))


@pytest.mark.parametrize("fmt", ["txt", "md", "csv", "docx", "xlsx", "pptx"])
def test_empty_text_and_office_do_not_count_markers_or_table_separators(tmp_path, fmt):
    path = tmp_path / f"empty.{fmt}"
    if fmt == "docx":
        from docx import Document
        doc = Document()
        doc.add_table(rows=2, cols=2)
        doc.save(path)
    elif fmt == "xlsx":
        from openpyxl import Workbook
        wb = Workbook()
        wb.active["B2"] = " "
        wb.save(path)
        wb.close()
    elif fmt == "pptx":
        from pptx import Presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(path)
    else:
        path.write_text(" \n\t", encoding="utf-8")
    diag = parse_source_with_diagnostics(path).diagnostics
    assert diag["extracted_chars"] == 0 and diag["empty_extraction"] is True


@pytest.mark.parametrize("text,excerpt,expected", [
    ("[[PAGE:1]]\nsame same\n[[PAGE:2]]\nunique", "same", {"status": "resolved", "locator": "PAGE:1"}),
    ("[[PAGE:1]]\nsame\n[[PAGE:2]]\nsame", "same", {"status": "ambiguous", "locators": ["PAGE:1", "PAGE:2"]}),
    ("[[PAGE:1]]\nＡＢＣ\\_d\n42", "ABC_d 42", {"status": "resolved", "locator": "PAGE:1"}),
    ("[[PAGE:1]]\nAlpha\n[[PAGE:2]]\nBeta", "Alpha Beta", {"status": "unresolved"}),
    ("[[PAGE:1]]\nAlpha", "alpha", {"status": "unresolved"}),
    ("[[PAGE:1]]\n[PAGE_PARSE_ERROR]", "PAGE_PARSE_ERROR", {"status": "unresolved"}),
    ("[[SLIDE:1]]\nText", "", {"status": "unresolved"}),
    ("[[TABLE:2]]\n[[TABLE:2:ROW:5]] value", "value", {"status": "resolved", "locator": "TABLE:2:ROW:5"}),
    ("[[PARA:17]] value", "value", {"status": "resolved", "locator": "PARA:17"}),
])
def test_deterministic_exact_locator(text, excerpt, expected):
    before = evidence_match(excerpt, text)
    assert resolve_evidence_locator(text, excerpt) == expected
    assert evidence_match(excerpt, text) == before


def test_locator_chunk_boundaries_and_oversized_fallback():
    units = [f"[[PAGE:{i}]]\n" + str(i) * 42 + "\n" for i in range(1, 4)]
    assert chunk_source_text("".join(units), 65) == units
    large = "[[PAGE:4]]\n" + "x" * 230 + "\nline two\n"
    text = "\n" + "".join(units) + large + "[[SLIDE:8]]\nend"
    chunks = chunk_source_text(text, 65)
    assert chunks == chunk_source_text(text, 65)
    assert "".join(chunks) == text
    assert all(0 < len(chunk) <= 65 for chunk in chunks)
    for marker in SOURCE_MARKER.finditer(text):
        assert sum(chunk.count(marker.group()) for chunk in chunks) == 1
    assert any("[[PAGE:4]]\nx" in chunk for chunk in chunks)


def test_long_mixed_locator_source_is_exactly_partitioned():
    text = "Preamble\n" + "".join(
        f"[[{kind}:{i}]]\n" + (f"unit {i} line\n" * (i % 9 + 1))
        for i in range(1, 201) for kind in ("PAGE", "SLIDE", "PARA", "TABLE")
    )
    for limit in (32, 75, 200, 500):
        chunks = chunk_source_text(text, limit)
        assert "".join(chunks) == text
        assert all(len(chunk) <= limit for chunk in chunks)
        assert [m.group() for chunk in chunks for m in SOURCE_MARKER.finditer(chunk)] == [m.group() for m in SOURCE_MARKER.finditer(text)]


def test_oversized_first_pdf_page_keeps_leading_separator_with_payload():
    text = "\n[[PAGE:1]]\n" + "body " * 100 + "\n[[PAGE:2]]\nlast page"
    chunks = chunk_source_text(text, 65)
    assert chunks[0].startswith("\n[[PAGE:1]]\nbody ")
    assert "".join(chunks) == text
    assert all(chunk.strip() and len(chunk) <= 65 for chunk in chunks)


@pytest.mark.parametrize("text", ["", "Plain text", "# Heading\n" + "paragraph\n" * 60])
def test_unmarked_chunking_is_unchanged(text):
    assert chunk_source_text(text, 50) == chunk_text(text, 50)


def test_impossible_chunk_limit_fails_without_splitting_marker():
    with pytest.raises(ValueError, match="locator marker"):
        chunk_source_text("[[SHEET:Capacity:ROW:31]] body", 10)
    with pytest.raises(ValueError, match="positive"):
        chunk_source_text("text", 0)


def test_unsupported_format(tmp_path):
    path = tmp_path / "unsupported.zip"
    path.write_bytes(b"unsupported")
    with pytest.raises(ParseError, match="Unsupported parser"):
        parse_source_with_diagnostics(path)
